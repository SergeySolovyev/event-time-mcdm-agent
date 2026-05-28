"""Build editable PPTX presentation from presentation.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Color palette --
DARK_NAVY = RGBColor(0x1a, 0x1a, 0x2e)
MID_NAVY = RGBColor(0x16, 0x21, 0x3e)
ACCENT_BLUE = RGBColor(0x0f, 0x34, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x1a, 0x1a, 0x2e)
TABLE_ALT_BG = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_NAVY, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_list,
                      font_size=18, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled runs per paragraph.
    runs_list: list of paragraphs, each paragraph is list of (text, bold, color) tuples.
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para_runs in enumerate(runs_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        for text, bold, color in para_runs:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Segoe UI"
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=17, color=DARK_NAVY):
    """Add a bulleted list. Each item can be (text,) or (bold_part, normal_part)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0

        if isinstance(item, tuple) and len(item) == 2:
            bold_part, normal_part = item
            r1 = p.add_run()
            r1.text = "  " + bold_part
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r1.font.name = "Segoe UI"
            r2 = p.add_run()
            r2.text = normal_part
            r2.font.size = Pt(font_size)
            r2.font.bold = False
            r2.font.color.rgb = color
            r2.font.name = "Segoe UI"
        else:
            r = p.add_run()
            r.text = "  " + str(item)
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return txBox


def add_table(slide, left, top, width, rows_data, col_widths=None,
              font_size=14):
    """Add a styled table. rows_data[0] = headers, rest = data rows."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_height = Inches(0.4 * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_NAVY

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return table_shape


def add_slide_number(slide, num, total=7):
    """Add slide number in bottom-right."""
    add_text_box(slide, 11.8, 7.0, 1.2, 0.4, f"{num} / {total}",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 -- Title
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, MID_NAVY)

add_text_box(slide1, 1.5, 1.8, 10.3, 1.2,
             "ERC-4626 Yield Vault\nwith MCDM Scoring",
             font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide1, 1.5, 3.5, 10.3, 0.6,
             "Autonomous DeFi Yield Optimization",
             font_size=24, color=RGBColor(0xAA, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide1, 1.5, 4.8, 10.3, 1.5,
             "Sergei Solovev\nsesesolovev@edu.hse.ru\nHSE University  |  April 2026\n\nProject 1: Custom DeFi Protocol",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 -- Why This Project?
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

add_text_box(slide2, 0.8, 0.4, 11.7, 0.7,
             "Why This Project?",
             font_size=36, bold=True, color=DARK_NAVY)

add_text_box(slide2, 0.8, 1.2, 11.7, 0.5,
             "The problem with DeFi yields",
             font_size=24, bold=True, color=ACCENT_BLUE)

add_text_box(slide2, 0.8, 1.8, 11.7, 0.8,
             "Lending protocols like Aave and Compound offer variable interest rates "
             "that change every block. Users who want the best yield face three challenges:",
             font_size=17, color=DARK_NAVY)

add_bullet_list(slide2, 0.8, 2.7, 11.7, 1.5, [
    ("Monitoring", " -- rates shift every ~12 seconds, impossible to track manually"),
    ("Gas costs", " -- each rebalance costs $2-50, eating into profits"),
    ("Risk", " -- chasing the highest APY ignores utilization risk and rate stability"),
], font_size=17)

add_text_box(slide2, 0.8, 4.4, 11.7, 0.5,
             "The idea: an autonomous agent",
             font_size=24, bold=True, color=ACCENT_BLUE)

add_text_box(slide2, 0.8, 5.0, 11.7, 1.2,
             "What if a software agent could watch the rates, evaluate multiple risk "
             "factors, and move funds automatically -- while proving every decision "
             "cryptographically on-chain?\n\n"
             "This is agentic DeFi: an off-chain agent that thinks, "
             "and an on-chain vault that verifies and executes.",
             font_size=17, color=DARK_NAVY)

add_slide_number(slide2, 2)

# ============================================================
# SLIDE 3 -- How It Works
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

add_text_box(slide3, 0.8, 0.4, 11.7, 0.7,
             "How It Works",
             font_size=36, bold=True, color=DARK_NAVY)

add_text_box(slide3, 0.8, 1.2, 11.7, 0.6,
             "User deposits USDC into the vault and receives aiUSDC shares.\n"
             "From that point, the agent manages everything automatically.",
             font_size=17, bold=True, color=DARK_NAVY)

# Architecture diagram as text box (monospace)
arch_text = (
    "+-------------------------------------------------------+\n"
    "|  Off-chain: Python Agent (runs every hour)             |\n"
    "|  1. Read APY, utilization, TVL, gas price              |\n"
    "|  2. Smooth rates with EMA (anti-manipulation)          |\n"
    "|  3. Score protocols with MCDM (4 factors)              |\n"
    "|  4. Sign decision with EIP-712 typed data              |\n"
    "+----------------------------+--------------------------+\n"
    "                             | signed tx\n"
    "+----------------------------+--------------------------+\n"
    "|  On-chain: AIVault.sol (ERC-4626 + UUPS proxy)        |\n"
    "|  - Verify keeper signature (ECDSA)                     |\n"
    "|  - Check nonce, timestamp, cooldown                    |\n"
    "|  - Execute rebalance via adapter                       |\n"
    "|  +-- AaveV3Adapter    +-- CompoundV3Adapter            |\n"
    "+-------------------------------------------------------+\n"
    "  Fallback: Chainlink Automation (if agent offline > 6h)"
)
add_text_box(slide3, 1.2, 2.1, 10.9, 4.0, arch_text,
             font_size=14, color=ACCENT_BLUE, font_name="Consolas")

add_text_box(slide3, 0.8, 6.2, 11.7, 0.8,
             "Key insight: The agent can be complex (multi-factor analysis), "
             "but the vault only trusts cryptographic proofs -- not the agent itself.",
             font_size=16, bold=True, color=DARK_NAVY)

add_slide_number(slide3, 3)

# ============================================================
# SLIDE 4 -- Key Formulas
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_text_box(slide4, 0.8, 0.4, 11.7, 0.7,
             "Key Formulas",
             font_size=36, bold=True, color=DARK_NAVY)

formulas = [
    [("ERC-4626 share price ", True, ACCENT_BLUE),
     ("(with inflation attack protection):", False, DARK_NAVY)],
    [("s = floor( a * (S + 10^6) / (A + 1) )", False, DARK_NAVY)],
    [],
    [("APY normalization ", True, ACCENT_BLUE),
     ("(cross-protocol, to annual 1e18 scale):", False, DARK_NAVY)],
    [("  Aave V3:  APY = liquidityRate_RAY / 10^9", False, DARK_NAVY)],
    [("  Compound V3:  APY = r_sec x 31,557,600", False, DARK_NAVY)],
    [],
    [("EMA smoothing ", True, ACCENT_BLUE),
     ("(dampens noise):", False, DARK_NAVY)],
    [("S_t = 0.3 * R_t + 0.7 * S_(t-1)", False, DARK_NAVY)],
    [],
    [("MCDM scoring model:", True, ACCENT_BLUE)],
    [("Score_i = 0.40*f_APY + 0.25*f_Risk + 0.20*f_Cost + 0.15*f_Stability",
      False, DARK_NAVY)],
    [],
    [("Rebalance if:  Score_best - Score_current >= 0.05", True, ACCENT_TEAL)],
]

add_rich_text_box(slide4, 0.8, 1.2, 11.7, 4.5, formulas, font_size=16)

add_text_box(slide4, 0.8, 5.8, 11.7, 1.2,
             "Example: Aave offers 6% APY but has 95% utilization (risky). "
             "Compound offers 5% but with 30% utilization (safe). "
             "Simple APY comparison picks Aave. Our MCDM model picks Compound "
             "-- the safer choice.",
             font_size=16, bold=False, color=DARK_NAVY)

add_slide_number(slide4, 4)

# ============================================================
# SLIDE 5 -- Security and Testing
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_text_box(slide5, 0.8, 0.4, 11.7, 0.7,
             "Security and Testing",
             font_size=36, bold=True, color=DARK_NAVY)

add_text_box(slide5, 0.8, 1.1, 5.5, 0.5,
             "7 threat mitigations:",
             font_size=20, bold=True, color=ACCENT_BLUE)

security_table = [
    ["Threat", "Protection"],
    ["Inflation attack", "Virtual shares (10^6 offset)"],
    ["Reentrancy", "ReentrancyGuard on all external functions"],
    ["Rate manipulation", "EMA + 5% jump guard"],
    ["Signature forgery", "EIP-712 domain + ECDSA verification"],
    ["Replay attack", "Sequential nonce + 5-min timestamp TTL"],
    ["Agent downtime", "Chainlink Automation fallback (6h)"],
    ["Rapid exploitation", "1-hour cooldown between rebalances"],
]

add_table(slide5, 0.8, 1.6, 5.8, security_table,
          col_widths=[2.0, 3.8], font_size=13)

add_text_box(slide5, 7.0, 1.1, 5.5, 0.5,
             "67 tests, 0 failures:",
             font_size=20, bold=True, color=ACCENT_BLUE)

test_table = [
    ["Category", "Tests", "Method"],
    ["Unit (Solidity)", "37", "Concrete + fuzz (1000 runs)"],
    ["Integration", "4", "Full lifecycle E2E"],
    ["Invariant", "6", "76,800+ random calls, 0 violations"],
    ["Python scoring", "20", "Pytest unit tests"],
]

add_table(slide5, 7.0, 1.6, 5.8, test_table,
          col_widths=[1.6, 0.8, 3.4], font_size=13)

add_slide_number(slide5, 5)

# ============================================================
# SLIDE 6 -- Deployment
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)

add_text_box(slide6, 0.8, 0.4, 11.7, 0.7,
             "Deployment",
             font_size=36, bold=True, color=DARK_NAVY)

add_text_box(slide6, 0.8, 1.2, 11.7, 0.5,
             "Ethereum Sepolia (Chain ID: 11155111). All contracts verified on Sourcify.",
             font_size=18, bold=True, color=DARK_NAVY)

contract_table = [
    ["Contract", "Address"],
    ["AaveV3Adapter", "0x8545D79f6FaB51EDc93Cf024fBD1FfAc98504ba1"],
    ["CompoundV3Adapter", "0xEB0D41F07691765314B9A45645Ee995d879c7ac7"],
    ["StrategyManager", "0x353469534dA4FB64d52Ae5059CEFd098557eBFa9"],
    ["AIVault (proxy)", "0x1324238b6F56Ccc785fC7f79Ca693546236Ad02C"],
]

add_table(slide6, 0.8, 1.9, 11.7, contract_table,
          col_widths=[3.0, 8.7], font_size=15)

add_text_box(slide6, 0.8, 4.2, 11.7, 0.5,
             "Tech stack:",
             font_size=20, bold=True, color=ACCENT_BLUE)

add_text_box(slide6, 0.8, 4.7, 11.7, 0.5,
             "Solidity 0.8.24, Python 3.12, Foundry, OpenZeppelin, Chainlink, Docker",
             font_size=17, color=DARK_NAVY)

add_text_box(slide6, 0.8, 5.4, 11.7, 0.5,
             "Design patterns used:",
             font_size=20, bold=True, color=ACCENT_BLUE)

add_bullet_list(slide6, 0.8, 5.9, 11.7, 1.2, [
    ("ERC-4626", " (tokenized vault standard)"),
    ("UUPS proxy", " (ERC-1967, upgradeable)"),
    ("Adapter / Strategy pattern", " (extensible to Morpho, Euler, etc.)"),
    ("EIP-712", " (typed data signing for agent decisions)"),
], font_size=17)

add_slide_number(slide6, 6)

# ============================================================
# SLIDE 7 -- Thank You
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, MID_NAVY)

add_text_box(slide7, 1.5, 2.2, 10.3, 1.0,
             "Thank You",
             font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide7, 1.5, 3.5, 10.3, 0.8,
             "Questions?",
             font_size=30, bold=True, color=RGBColor(0xAA, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide7, 1.5, 5.0, 10.3, 1.5,
             "GitHub: github.com/SergeySolovyev/ai-yield-vault\n"
             "Contact: sesesolovev@edu.hse.ru | @Sergey_Solovjov | www.sergeisolovev.com",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
out_path = r"D:\DeFi\DeFi-Vega Project\docs\presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
